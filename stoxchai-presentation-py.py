#!/usr/bin/env python
# stoxchai_presentation.py - Create a PowerPoint presentation for StoxChai

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import os
import base64
from PIL import Image

# Create presentation
prs = Presentation()

# Define StoxChai colors
STOXCHAI_BLUE = RGBColor(30, 50, 140)  # Dark blue
STOXCHAI_LIGHT_BLUE = RGBColor(69, 123, 157)  # Light blue
STOXCHAI_ACCENT = RGBColor(239, 71, 111)  # Pink accent
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 102, 102)
BLACK = RGBColor(0, 0, 0)
GREEN = RGBColor(5, 150, 105)
RED = RGBColor(220, 38, 38)

# Set slide dimensions (16:9 ratio)
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# def add_footer(slide, text="StoxChai | Financial Analysis with AI | May 2025"):
#     """Add footer to slide"""
#     footer = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.4))
#     footer_frame = footer.text_frame
#     footer_frame.text = text
#     footer_frame.paragraphs[0].font.size = Pt(10)
#     footer_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE
#     footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Function to load the image from assets
def save_images():
    """Save the screenshots as image files to be used in the presentation"""
    # Define the directory
    if not os.path.exists("assets"):
        os.makedirs("assets")
    
    # The image paths
    image_paths = [
        "assets/ai_analysis.png",
        "assets/chat_interface.png",
        "assets/price_chart.png",
        "assets/news_sentiment.png"
    ]
    
    # Check if images already exist
    if all(os.path.exists(path) for path in image_paths):
        return image_paths
    
    # Example message - in a real scenario, you would need to have these images
    print("Please ensure the following images exist in the 'assets' folder:")
    for path in image_paths:
        print(f"- {path}")
    
    # Return the paths anyway
    return image_paths

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Add gradient background shape
bg_shape = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
)
bg_shape.fill.gradient()
bg_shape.fill.gradient_angle = 45
bg_shape.fill.gradient_stops[0].color.rgb = STOXCHAI_BLUE
bg_shape.fill.gradient_stops[1].color.rgb = STOXCHAI_LIGHT_BLUE
bg_shape.line.fill.background()

# Add title
title = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
title_frame = title.text_frame
title_frame.text = "StoxChai"
title_frame.paragraphs[0].font.size = Pt(60)
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add subtitle
subtitle = slide1.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
subtitle_frame = subtitle.text_frame
subtitle_frame.text = "Advanced Stock Analysis with AI"
subtitle_frame.paragraphs[0].font.size = Pt(32)
subtitle_frame.paragraphs[0].font.color.rgb = WHITE
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# add_footer(slide1)

# Slide 2: Overview
slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
slide2.shapes.title.text = "Comprehensive Stock Analysis Platform"
slide2.shapes.title.text_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE

# Add content
content = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
content_frame = content.text_frame

p1 = content_frame.add_paragraph()
p1.text = "StoxChai is an all-in-one analysis platform that helps investors make data-driven decisions by providing:"
p1.font.size = Pt(18)
p1.space_after = Pt(12)

features = [
    "Real-time stock data with interactive price charts",
    "Technical analysis with moving averages and trend detection",
    "News sentiment analysis to gauge market perception",
    "Volume analysis to identify unusual trading activity",
    "AI-powered insights through a conversational interface",
    "Retrieval Augmented Generation (RAG) for context-aware responses"
]

for feature in features:
    p = content_frame.add_paragraph()
    p.text = f"• {feature}"
    p.font.size = Pt(18)
    p.level = 1

# add_footer(slide2)

# Slide 3: Key Functionalities
slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
slide3.shapes.title.text = "Key Functionalities"
slide3.shapes.title.text_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE

# Create a card layout for functionalities
functionalities = [
    {
        "title": "Price Analysis",
        "description": "Interactive candlestick charts with moving averages and technical indicators.",
        "icon": "📈"
    },
    {
        "title": "Volume Analysis",
        "description": "Track trading volume patterns and identify significant spikes.",
        "icon": "📊"
    },
    {
        "title": "News Sentiment",
        "description": "NLP-driven analysis of news articles with NLTK sentiment scoring.",
        "icon": "📰"
    },
    {
        "title": "AI Assistant",
        "description": "RAG-powered insights with vector similarity search using context retrieval.",
        "icon": "🤖"
    }
]

y_pos = 1.5
for i, func in enumerate(functionalities):
    # Create function card
    card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(240, 240, 250)  # Light blue background
    
    # Add icon text
    icon = slide3.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.15), Inches(0.5), Inches(0.5))
    icon_frame = icon.text_frame
    icon_frame.text = func["icon"]
    icon_frame.paragraphs[0].font.size = Pt(24)
    
    # Add title
    title_box = slide3.shapes.add_textbox(Inches(1.4), Inches(y_pos + 0.1), Inches(7.8), Inches(0.3))
    title_frame = title_box.text_frame
    title_frame.text = func["title"]
    title_frame.paragraphs[0].font.size = Pt(18)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE
    
    # Add description
    desc_box = slide3.shapes.add_textbox(Inches(1.4), Inches(y_pos + 0.4), Inches(7.8), Inches(0.3))
    desc_frame = desc_box.text_frame
    desc_frame.text = func["description"]
    desc_frame.paragraphs[0].font.size = Pt(14)
    desc_frame.paragraphs[0].font.color.rgb = GRAY
    
    y_pos += 1

# add_footer(slide3)

# Save the image paths
image_paths = save_images()

# Slide 4: AI Analysis Screenshot
slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
slide4.shapes.title.text = "AI Analysis"
slide4.shapes.title.text_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE

# Add description
desc = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
desc_frame = desc.text_frame
desc_frame.text = "Generate comprehensive stock analyses with key metrics, performance trends, and sector comparisons."
desc_frame.paragraphs[0].font.size = Pt(16)
desc_frame.paragraphs[0].font.color.rgb = GRAY

# Add screenshot placeholder
try:
    if os.path.exists(image_paths[0]):
        slide4.shapes.add_picture(image_paths[0], Inches(0.5), Inches(1.8), width=Inches(9))
    else:
        img_placeholder = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(9), Inches(3)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
        
        # Add text to indicate missing image
        text_box = slide4.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(0.5))
        text_box.text_frame.text = "AI Analysis Screenshot"
        text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
except Exception as e:
    print(f"Error adding image: {e}")

# add_footer(slide4)

# Slide 5: Chat Interface Screenshot
slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
slide5.shapes.title.text = "AI Assistant Chat"
slide5.shapes.title.text_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE

# Add description
desc = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
desc_frame = desc.text_frame
desc_frame.text = "Interact with an AI assistant that provides context-aware responses about stocks using RAG technology."
desc_frame.paragraphs[0].font.size = Pt(16)
desc_frame.paragraphs[0].font.color.rgb = GRAY

# Add screenshot placeholder
try:
    if os.path.exists(image_paths[1]):
        slide5.shapes.add_picture(image_paths[1], Inches(0.5), Inches(1.8), width=Inches(9))
    else:
        img_placeholder = slide5.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(9), Inches(3)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
        
        # Add text to indicate missing image
        text_box = slide5.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(0.5))
        text_box.text_frame.text = "Chat Interface Screenshot"
        text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
except Exception as e:
    print(f"Error adding image: {e}")

# add_footer(slide5)

# Slide 6: Price Chart Screenshot
slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
slide6.shapes.title.text = "Interactive Price Charts"
slide6.shapes.title.text_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE

# Add description
desc = slide6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
desc_frame = desc.text_frame
desc_frame.text = "Analyze stock price movements with candlestick charts, multiple timeframes, and technical indicators."
desc_frame.paragraphs[0].font.size = Pt(16)
desc_frame.paragraphs[0].font.color.rgb = GRAY

# Add screenshot placeholder
try:
    if os.path.exists(image_paths[2]):
        slide6.shapes.add_picture(image_paths[2], Inches(0.5), Inches(1.8), width=Inches(9))
    else:
        img_placeholder = slide6.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(9), Inches(3)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
        
        # Add text to indicate missing image
        text_box = slide6.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(0.5))
        text_box.text_frame.text = "Price Chart Screenshot"
        text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
except Exception as e:
    print(f"Error adding image: {e}")

# add_footer(slide6)

# Slide 7: News Sentiment Screenshot
slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
slide7.shapes.title.text = "News Sentiment Analysis"
slide7.shapes.title.text_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE

# Add description
desc = slide7.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
desc_frame = desc.text_frame
desc_frame.text = "Track market perception through automated analysis of news articles with sentiment scoring visualization."
desc_frame.paragraphs[0].font.size = Pt(16)
desc_frame.paragraphs[0].font.color.rgb = GRAY

# Add screenshot placeholder
try:
    if os.path.exists(image_paths[3]):
        slide7.shapes.add_picture(image_paths[3], Inches(0.5), Inches(1.8), width=Inches(9))
    else:
        img_placeholder = slide7.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(9), Inches(3)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
        
        # Add text to indicate missing image
        text_box = slide7.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(0.5))
        text_box.text_frame.text = "News Sentiment Screenshot"
        text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
except Exception as e:
    print(f"Error adding image: {e}")

# add_footer(slide7)

# Slide 8: Target Users
slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
slide8.shapes.title.text = "Who Benefits From StoxChai?"
slide8.shapes.title.text_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE

# Create a card layout for user types
users = [
    {
        "title": "Day Traders",
        "description": "Real-time technical indicators for quick decision making",
        "icon": "⚡"
    },
    {
        "title": "Long-Term Investors",
        "description": "Fundamental data and trend analysis for portfolio decisions",
        "icon": "📝"
    },
    {
        "title": "Financial Advisors",
        "description": "Client-ready visualizations and summaries",
        "icon": "👔"
    },
    {
        "title": "Finance Students",
        "description": "Learning platform for market analysis techniques",
        "icon": "🎓"
    }
]

y_pos = 1.5
for i, user in enumerate(users):
    # Create user card
    card = slide8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(240, 240, 250)  # Light blue background
    
    # Add icon text
    icon = slide8.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.15), Inches(0.5), Inches(0.5))
    icon_frame = icon.text_frame
    icon_frame.text = user["icon"]
    icon_frame.paragraphs[0].font.size = Pt(24)
    
    # Add title
    title_box = slide8.shapes.add_textbox(Inches(1.4), Inches(y_pos + 0.1), Inches(7.8), Inches(0.3))
    title_frame = title_box.text_frame
    title_frame.text = user["title"]
    title_frame.paragraphs[0].font.size = Pt(18)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE
    
    # Add description
    desc_box = slide8.shapes.add_textbox(Inches(1.4), Inches(y_pos + 0.4), Inches(7.8), Inches(0.3))
    desc_frame = desc_box.text_frame
    desc_frame.text = user["description"]
    desc_frame.paragraphs[0].font.size = Pt(14)
    desc_frame.paragraphs[0].font.color.rgb = GRAY
    
    y_pos += 1

# add_footer(slide8)

# Slide 9: Technology Stack
slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
slide9.shapes.title.text = "Technology Stack"
slide9.shapes.title.text_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE

# Create columns for frontend and backend
col_left = slide9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.3), Inches(3.3)
)
col_left.fill.solid()
col_left.fill.fore_color.rgb = RGBColor(240, 240, 250)

col_right = slide9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.3), Inches(3.3)
)
col_right.fill.solid()
col_right.fill.fore_color.rgb = RGBColor(240, 240, 250)

# Add titles to columns
left_title = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.4))
left_title_frame = left_title.text_frame
left_title_frame.text = "Frontend"
left_title_frame.paragraphs[0].font.size = Pt(18)
left_title_frame.paragraphs[0].font.bold = True
left_title_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE
left_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

right_title = slide9.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.4))
right_title_frame = right_title.text_frame
right_title_frame.text = "Backend & AI"
right_title_frame.paragraphs[0].font.size = Pt(18)
right_title_frame.paragraphs[0].font.bold = True
right_title_frame.paragraphs[0].font.color.rgb = STOXCHAI_BLUE
right_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add frontend technologies
frontend_tech = [
    "Streamlit - Python-based UI framework",
    "Plotly - Interactive data visualization",
    "Custom CSS - Responsive design",
    "Interactive components - Filters, charts, inputs"
]

frontend_list = slide9.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(4), Inches(2.5))
frontend_frame = frontend_list.text_frame

for tech in frontend_tech:
    p = frontend_frame.add_paragraph()
    p.text = f"• {tech}"
    p.font.size = Pt(14)
    p.space_after = Pt(12)

# Add backend technologies
backend_tech = [
    "Python - Core application logic",
    "yfinance - Real-time market data",
    "Pandas/NumPy - Data processing",
    "Ollama - Local LLM integration",
    "Custom RAG system - Context-aware AI",
    "NLTK - Sentiment analysis"
]

backend_list = slide9.shapes.add_textbox(Inches(5.4), Inches(2.0), Inches(4), Inches(2.5))
backend_frame = backend_list.text_frame

for tech in backend_tech:
    p = backend_frame.add_paragraph()
    p.text = f"• {tech}"
    p.font.size = Pt(14)
    p.space_after = Pt(12)

# add_footer(slide9)

# Slide 10: Thank You
slide10 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Add gradient background
bg_shape = slide10.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
)
bg_shape.fill.gradient()
bg_shape.fill.gradient_angle = 45
bg_shape.fill.gradient_stops[0].color.rgb = STOXCHAI_BLUE
bg_shape.fill.gradient_stops[1].color.rgb = STOXCHAI_LIGHT_BLUE
bg_shape.line.fill.background()

# Add title
title = slide10.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
title_frame = title.text_frame
title_frame.text = "Thank You!"
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add contact info
subtitle = slide10.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(0.5))
subtitle_frame = subtitle.text_frame
subtitle_frame.text = "Advanced Stock Analysis with AI"
subtitle_frame.paragraphs[0].font.size = Pt(24)
subtitle_frame.paragraphs[0].font.color.rgb = WHITE
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add link or QR code
contact = slide10.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
contact_frame = contact.text_frame
contact_frame.text = "github.com/yourusername/stoxchai"
contact_frame.paragraphs[0].font.size = Pt(18)
contact_frame.paragraphs[0].font.color.rgb = WHITE
contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# add_footer(slide10, "StoxChai | Financial Analysis with AI | © 2025")

# Save the presentation
prs.save('StoxChai_Presentation.pptx')
print("PowerPoint presentation 'StoxChai_Presentation.pptx' has been created successfully!")
